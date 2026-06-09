"""Export cost analysis data to JSON, CSV, and Markdown formats."""
from __future__ import annotations

import json
import csv
import io
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

import pandas as pd


def _timestamp() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")


def _format_cost_md(amount: float) -> str:
    if abs(amount) >= 1_000_000:
        return f"${amount / 1_000_000:,.1f}M"
    elif abs(amount) >= 10_000:
        return f"${amount / 1_000:,.1f}K"
    elif abs(amount) >= 1_000:
        return f"${amount:,.0f}"
    return f"${amount:,.2f}"


class CostExporter:
    """Export cost analysis results to multiple formats."""

    def __init__(
        self,
        data: dict[str, pd.DataFrame],
        days: int,
        anomalies: Optional[pd.DataFrame] = None,
        insights: Optional[dict] = None,
        story: Optional[str] = None,
        redact: bool = True,
    ):
        self.data = data
        self.days = days
        self.anomalies = anomalies if anomalies is not None else pd.DataFrame()
        self.insights = insights or {}
        self.story = story or ""
        self.ts = _timestamp()
        self.redact = redact

    def _redact_df(self, df: pd.DataFrame) -> pd.DataFrame:
        """Apply redaction to DataFrame columns that may contain PII."""
        if not self.redact or df.empty:
            return df
        from kulshan.redact import redact_account_id, redact_text
        result = df.copy()
        # Redact account columns
        for col in result.columns:
            if col in ("account", "account_id", "accountId"):
                result[col] = result[col].apply(
                    lambda x: redact_account_id(str(x)) if pd.notna(x) else x
                )
            elif col in ("title", "description", "recommended_action", "why_it_matters"):
                result[col] = result[col].apply(
                    lambda x: redact_text(str(x)) if pd.notna(x) else x
                )
        return result

    def _redact_payload(self, payload: dict) -> dict:
        """Apply redaction to a JSON payload dict."""
        if not self.redact:
            return payload
        from kulshan.redact import redact_payload
        return redact_payload(payload)

    def _base_name(self, prefix: str = "cost_report") -> str:
        return f"{prefix}_{self.ts}"

    # ── JSON ──────────────────────────────────────────────────────────────

    def export_json(self, path: Optional[str] = None) -> str:
        path = path or f"{self._base_name()}.json"
        payload = {
            "generated": datetime.now(timezone.utc).isoformat(),
            "days_analyzed": self.days,
            "story": self.story,
            "insights": self.insights,
        }
        for key, df in self.data.items():
            if not df.empty:
                records = df.copy()
                # Convert dates to strings for JSON serialization
                for col in records.select_dtypes(include=["datetime64"]).columns:
                    records[col] = records[col].dt.strftime("%Y-%m-%d")
                payload[key] = records.to_dict(orient="records")

        if not self.anomalies.empty:
            anom = self.anomalies.copy()
            for col in anom.select_dtypes(include=["datetime64"]).columns:
                anom[col] = anom[col].dt.strftime("%Y-%m-%d")
            payload["anomalies"] = anom.to_dict(orient="records")

        Path(path).write_text(json.dumps(self._redact_payload(payload), indent=2, default=str))
        return path

    # ── CSV ───────────────────────────────────────────────────────────────

    def export_csv(self, path: Optional[str] = None) -> str:
        path = path or f"{self._base_name()}.csv"
        # Combine all dimension data into one CSV with a "dimension" column
        frames = []
        for key, df in self.data.items():
            if df.empty:
                continue
            tmp = df.copy()
            tmp["dimension"] = key
            frames.append(tmp)

        if frames:
            combined = pd.concat(frames, ignore_index=True)
            self._redact_df(combined).to_csv(path, index=False)
        else:
            Path(path).write_text("")
        return path

    # ── Markdown ──────────────────────────────────────────────────────────

    def export_markdown(self, path: Optional[str] = None) -> str:
        path = path or f"{self._base_name()}.md"
        lines = [
            f"# Cost Analysis Report",
            f"",
            f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')} · Last {self.days} days",
            f"",
        ]

        if self.story:
            lines += [f"## Cost Story", f"", self.story, f""]

        # Summary
        svc_df = self.data.get("service", pd.DataFrame())
        if not svc_df.empty:
            total = svc_df["cost"].sum()
            daily_avg = total / max(self.days, 1)
            num_svc = svc_df[svc_df.columns[1]].nunique() if len(svc_df.columns) > 1 else 0
            lines += [
                f"## Summary",
                f"",
                f"| Metric | Value |",
                f"|--------|-------|",
                f"| Total Spend | {_format_cost_md(total)} |",
                f"| Daily Average | {_format_cost_md(daily_avg)} |",
                f"| Active Services | {num_svc} |",
                f"",
            ]

        # Cost tables per dimension
        for key, df in self.data.items():
            if df.empty:
                continue
            totals = df.groupby(key)["cost"].sum().sort_values(ascending=False).head(9)
            grand = totals.sum()
            lines += [
                f"## Cost by {key.replace('_', ' ').title()}",
                f"",
                f"| {key.title()} | Cost | % Total |",
                f"|---|---:|---:|",
            ]
            for name, cost in totals.items():
                pct = cost / grand * 100 if grand > 0 else 0
                lines.append(f"| {name} | {_format_cost_md(cost)} | {pct:.1f}% |")
            lines += [f"| **Total** | **{_format_cost_md(grand)}** | 100% |", f""]

        # Anomalies
        if not self.anomalies.empty:
            lines += [f"## Anomalies ({len(self.anomalies)} detected)", f""]
            group_col = self.anomalies.columns[0]
            lines += [
                f"| {group_col.title()} | Cost | Avg | Deviation | Score | Methods | Severity |",
                f"|---|---:|---:|---:|---:|---|---|",
            ]
            for _, row in self.anomalies.head(9).iterrows():
                lines.append(
                    f"| {row[group_col]} | ${row['latest_cost']:,.2f} | "
                    f"${row['avg_cost']:,.2f} | {row.get('pct_change', 0):+.1f}% | "
                    f"{row.get('score', 0):.1f}σ | {row.get('methods', '')} | "
                    f"{row.get('severity', 'info')} |"
                )
            lines.append(f"")

        Path(path).write_text("\n".join(lines), encoding="utf-8")
        return path

    # ── PDF ──────────────────────────────────────────────────────────────

    def export_pdf(self, path: Optional[str] = None, html_content: Optional[str] = None) -> str:
        """Export report as PDF. Requires weasyprint (optional dep)."""
        path = path or f"{self._base_name()}.pdf"
        try:
            from weasyprint import HTML as WeasyHTML
        except ImportError:
            raise RuntimeError(
                "PDF export requires weasyprint. Install with: pip install weasyprint"
            )
        if html_content:
            WeasyHTML(string=html_content).write_pdf(path)
        else:
            # Generate a temporary HTML and convert
            from .html_report import generate_html_report
            tmp_html = f"{self._base_name()}_tmp.html"
            generate_html_report(self.data, tmp_html, self.days)
            WeasyHTML(filename=tmp_html).write_pdf(path)
            Path(tmp_html).unlink(missing_ok=True)
        return path

    # ── Excel (.xlsx) with charts ────────────────────────────────────────

    def export_excel(self, path: Optional[str] = None) -> str:
        """Export to Excel with embedded charts. Requires openpyxl."""
        path = path or f"{self._base_name()}.xlsx"
        try:
            from openpyxl import Workbook
            from openpyxl.chart import BarChart, Reference, LineChart
        except ImportError:
            raise RuntimeError("Excel export requires openpyxl. Install with: pip install openpyxl")

        wb = Workbook()

        # Summary sheet
        ws = wb.active
        ws.title = "Summary"
        ws.append(["Metric", "Value"])
        svc_df = self.data.get("service", pd.DataFrame())
        if not svc_df.empty:
            total = svc_df["cost"].sum()
            ws.append(["Total Spend", round(total, 2)])
            ws.append(["Daily Average", round(total / max(self.days, 1), 2)])
            ws.append(["Days Analyzed", self.days])
            ws.append(["Generated", datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")])
            if self.story:
                ws.append([])
                ws.append(["Cost Story", self.story])

        # Service breakdown with chart
        for key in ["service", "account", "region"]:
            if key not in self.data or self.data[key].empty:
                continue
            df = self.data[key]
            totals = df.groupby(key)["cost"].sum().sort_values(ascending=False).head(9)
            ws_dim = wb.create_sheet(title=key.title()[:31])
            ws_dim.append([key.title(), "Cost (USD)", "% of Total"])
            grand = totals.sum()
            for name, cost in totals.items():
                pct = cost / grand * 100 if grand > 0 else 0
                ws_dim.append([str(name), round(cost, 2), round(pct, 1)])

            # Add bar chart
            if len(totals) > 1:
                chart = BarChart()
                chart.title = f"Cost by {key.title()}"
                chart.y_axis.title = "Cost (USD)"
                chart.style = 10
                data_ref = Reference(ws_dim, min_col=2, min_row=1, max_row=len(totals) + 1)
                cats_ref = Reference(ws_dim, min_col=1, min_row=2, max_row=len(totals) + 1)
                chart.add_data(data_ref, titles_from_data=True)
                chart.set_categories(cats_ref)
                chart.width = 20
                chart.height = 12
                ws_dim.add_chart(chart, f"E2")

        # Daily trend with line chart
        if not svc_df.empty:
            daily = svc_df.groupby("date")["cost"].sum().reset_index().sort_values("date")
            if len(daily) > 1:
                ws_trend = wb.create_sheet(title="Daily Trend")
                ws_trend.append(["Date", "Cost"])
                for _, row in daily.iterrows():
                    ws_trend.append([row["date"].strftime("%Y-%m-%d"), round(row["cost"], 2)])
                chart = LineChart()
                chart.title = "Daily Cost Trend"
                chart.y_axis.title = "Cost (USD)"
                chart.style = 10
                data_ref = Reference(ws_trend, min_col=2, min_row=1, max_row=len(daily) + 1)
                cats_ref = Reference(ws_trend, min_col=1, min_row=2, max_row=len(daily) + 1)
                chart.add_data(data_ref, titles_from_data=True)
                chart.set_categories(cats_ref)
                chart.width = 25
                chart.height = 12
                ws_trend.add_chart(chart, "D2")

        # Anomalies sheet
        if not self.anomalies.empty:
            ws_anom = wb.create_sheet(title="Anomalies")
            cols = [c for c in self.anomalies.columns if c != "date"]
            ws_anom.append(cols)
            for _, row in self.anomalies.head(20).iterrows():
                ws_anom.append([str(row.get(c, "")) for c in cols])

        wb.save(path)
        return path

    # ── PowerPoint (.pptx) ────────────────────────────────────────────

    def export_pptx(self, path: Optional[str] = None) -> str:
        """Export to PowerPoint. Requires python-pptx."""
        path = path or f"{self._base_name()}.pptx"
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise RuntimeError("PowerPoint export requires python-pptx. Install with: pip install python-pptx")

        prs = Presentation()

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Cost Analysis Report"
        slide.placeholders[1].text = f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d')} · Last {self.days} days"

        # Cost Story slide
        if self.story:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Executive Summary"
            slide.placeholders[1].text = self.story

        # Service breakdown slide
        svc_df = self.data.get("service", pd.DataFrame())
        if not svc_df.empty:
            totals = svc_df.groupby("service")["cost"].sum().sort_values(ascending=False).head(9)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Top Services by Cost"
            body = slide.placeholders[1]
            body.text = ""
            for name, cost in totals.items():
                body.text += f"• {name}: {_format_cost_md(cost)}\n"

        # Anomalies slide
        if not self.anomalies.empty:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Anomalies ({len(self.anomalies)} detected)"
            body = slide.placeholders[1]
            body.text = ""
            group_col = self.anomalies.columns[0]
            for _, row in self.anomalies.head(5).iterrows():
                body.text += (f"• {row[group_col]}: {_format_cost_md(row['latest_cost'])} "
                              f"({row.get('pct_change', 0):+.1f}% vs avg)\n")

        prs.save(path)
        return path

    # ── Slack webhook ─────────────────────────────────────────────────

    def send_slack(self, webhook_url: str) -> bool:
        """Explicitly send a cost summary to the user-provided Slack webhook."""
        import urllib.request
        import json as _json

        svc_df = self.data.get("service", pd.DataFrame())
        total = svc_df["cost"].sum() if not svc_df.empty else 0

        blocks = [
            {"type": "header", "text": {"type": "plain_text", "text": "📊 Cost Analysis Report"}},
            {"type": "section", "text": {"type": "mrkdwn", "text": (
                f"*Period:* Last {self.days} days\n"
                f"*Total Spend:* {_format_cost_md(total)}\n"
                f"*Daily Average:* {_format_cost_md(total / max(self.days, 1))}"
            )}},
        ]

        if self.story:
            blocks.append({"type": "section", "text": {"type": "mrkdwn", "text": f"_{self.story}_"}})

        if not self.anomalies.empty:
            anom_count = len(self.anomalies)
            blocks.append({"type": "section", "text": {"type": "mrkdwn", "text": f"⚠️ *{anom_count} anomalies detected*"}})

        payload = _json.dumps({"blocks": blocks}).encode("utf-8")
        req = urllib.request.Request(webhook_url, data=payload, headers={"Content-Type": "application/json"})
        try:
            urllib.request.urlopen(req)
            return True
        except Exception:
            return False



    # ── Export all ─────────────────────────────────────────────────────────

    def export_all(self, html_path: Optional[str] = None) -> list[str]:
        """Export to all formats. Returns list of created file paths."""
        paths = [
            self.export_json(),
            self.export_csv(),
            self.export_markdown(),
        ]
        try:
            paths.append(self.export_pdf())
        except RuntimeError:
            pass
        try:
            paths.append(self.export_excel())
        except RuntimeError:
            pass
        try:
            paths.append(self.export_pptx())
        except RuntimeError:
            pass
        return paths


    # ── iCal export (.ics) ────────────────────────────────────────────

    def export_ical(self, path: Optional[str] = None) -> str:
        """Export anomalies and forecast milestones as iCal events."""
        path = path or f"{self._base_name()}.ics"
        events = []

        # Anomaly events
        if not self.anomalies.empty:
            for _, row in self.anomalies.head(9).iterrows():
                group_col = self.anomalies.columns[0]
                svc = str(row[group_col])
                cost = row.get("latest_cost", 0)
                pct = row.get("pct_change", 0)
                sev = row.get("severity", "info")
                date_val = row.get("date", "")
                if hasattr(date_val, "strftime"):
                    dt = date_val.strftime("%Y%m%d")
                else:
                    dt = str(date_val).replace("-", "")[:8]

                events.append(
                    f"BEGIN:VEVENT\n"
                    f"DTSTART;VALUE=DATE:{dt}\n"
                    f"SUMMARY:⚠️ AWS Anomaly: {svc} ({pct:+.1f}%)\n"
                    f"DESCRIPTION:Cost: ${cost:,.2f} | Severity: {sev} | {row.get('methods', '')}\n"
                    f"END:VEVENT"
                )

        # Summary event
        svc_df = self.data.get("service", pd.DataFrame())
        if not svc_df.empty:
            total = svc_df["cost"].sum()
            today = datetime.now(timezone.utc).strftime("%Y%m%d")
            events.append(
                f"BEGIN:VEVENT\n"
                f"DTSTART;VALUE=DATE:{today}\n"
                f"SUMMARY:📊 AWS Cost Report: {_format_cost_md(total)} ({self.days}d)\n"
                f"DESCRIPTION:{self.story}\n"
                f"END:VEVENT"
            )

        ical = (
            "BEGIN:VCALENDAR\n"
            "VERSION:2.0\n"
            "PRODID:-//Mission FinOps//EN\n"
            + "\n".join(events) + "\n"
            "END:VCALENDAR"
        )
        Path(path).write_text(ical, encoding="utf-8")
        return path
